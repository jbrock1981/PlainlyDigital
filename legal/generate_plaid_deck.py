#!/usr/bin/env python3
"""
Generate the Patet × Plaid partnership / IAM-context slide deck.

Audience: Plaid's IAM-review / partnership team. The deck is positioning
context, not a fundraising pitch — it answers "what is Patet, why does
the Plaid integration matter to your product, and are you a stable
partner?" in 12 slides.

Sources for content:
  - Competitor research (general-purpose agent, 2026-05-17, validated
    via WebSearch — see PATET_PLAID_PARTNERSHIP_DECK_NOTES.md)
  - Patet feature audit (Explore agent, 2026-05-17, against the codebase)
  - InfoSec Policy, Access Control Policy, Data Retention Policy
  - CLAUDE.md (canonical source for product architecture)

Run:
    /tmp/docxenv/bin/python legal/generate_plaid_deck.py

Output:
    legal/PATET_PLAID_PARTNERSHIP_DECK.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "PATET_PLAID_PARTNERSHIP_DECK.pptx"

# ─── Brand colors ─────────────────────────────────────────────────────────────
INK = RGBColor(0x0F, 0x0F, 0x0F)        # near-black, primary text on light
PAPER = RGBColor(0xFF, 0xFF, 0xFF)      # white, body bg
CHARCOAL = RGBColor(0x1C, 0x1C, 0x1E)   # title-bar dark
GREEN = RGBColor(0x10, 0xB9, 0x81)      # Patet brand green
MUTED = RGBColor(0x6B, 0x72, 0x80)      # secondary text
RULE = RGBColor(0xE5, 0xE7, 0xEB)       # hairline divider
GREEN_TINT = RGBColor(0xEC, 0xFD, 0xF5)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_TINT = RGBColor(0xFE, 0xF3, 0xC7)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


# ─── Helpers ──────────────────────────────────────────────────────────────────


def _set_solid(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, left, top, width, height, color):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _set_solid(rect, color)
    rect.shadow.inherit = False
    return rect


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    *,
    size=14,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font="Inter",
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(
    slide,
    paragraphs,
    left,
    top,
    width,
    height,
    *,
    size=14,
    color=INK,
    line_spacing=1.25,
    font="Inter",
    align=PP_ALIGN.LEFT,
):
    """paragraphs = list of (text, bold) tuples. First paragraph reuses the
    default; subsequent ones get add_paragraph()."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, (text, bold) in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, bullets, left, top, width, height, *, size=14, color=INK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.35
        if i > 0:
            p.space_before = Pt(6)
        # bullet via Unicode middot
        dot = p.add_run()
        dot.text = "•  "
        dot.font.name = "Inter"
        dot.font.size = Pt(size)
        dot.font.color.rgb = GREEN
        dot.font.bold = True
        body = p.add_run()
        body.text = b
        body.font.name = "Inter"
        body.font.size = Pt(size)
        body.font.color.rgb = color
    return tb


def slide_header(slide, eyebrow, title, *, eyebrow_color=GREEN):
    """Top eyebrow + big title. Returns the bottom Y for content placement."""
    add_text(
        slide,
        eyebrow.upper(),
        MARGIN,
        Inches(0.5),
        SLIDE_W - 2 * MARGIN,
        Inches(0.3),
        size=11,
        bold=True,
        color=eyebrow_color,
    )
    add_text(
        slide,
        title,
        MARGIN,
        Inches(0.85),
        SLIDE_W - 2 * MARGIN,
        Inches(0.9),
        size=32,
        bold=True,
        color=INK,
    )
    # hairline rule
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.8), SLIDE_W - 2 * MARGIN, Emu(9525)
    )
    _set_solid(rule, RULE)
    return Inches(2.0)


def slide_footer(slide, slide_num, total):
    add_text(
        slide,
        "Patet · Plainly Digital LLC",
        MARGIN,
        Inches(7.05),
        Inches(5),
        Inches(0.3),
        size=9,
        color=MUTED,
    )
    add_text(
        slide,
        f"{slide_num} / {total}",
        SLIDE_W - MARGIN - Inches(1),
        Inches(7.05),
        Inches(1),
        Inches(0.3),
        size=9,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


# ─── Slides ───────────────────────────────────────────────────────────────────

TOTAL_SLIDES = 12


def slide_01_title():
    s = prs.slides.add_slide(BLANK)
    # full-bleed charcoal panel
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, CHARCOAL)
    # accent block
    add_rect(s, Inches(0.6), Inches(2.6), Inches(0.18), Inches(1.6), GREEN)
    # logotype
    add_text(
        s,
        "PATET",
        Inches(1.0),
        Inches(2.55),
        Inches(8),
        Inches(0.7),
        size=42,
        bold=True,
        color=PAPER,
    )
    # tagline
    add_text(
        s,
        "Money skills + your real bank data + an AI coach who knows what you’re actually spending on.",
        Inches(1.0),
        Inches(3.35),
        Inches(11.3),
        Inches(0.9),
        size=20,
        color=PAPER,
    )
    # context line
    add_text(
        s,
        "Plaid IAM partnership briefing  ·  prepared May 2026",
        Inches(1.0),
        Inches(4.55),
        Inches(11.3),
        Inches(0.4),
        size=13,
        color=RGBColor(0xC0, 0xC0, 0xC0),
    )
    # plainly digital footer
    add_text(
        s,
        "Plainly Digital LLC  ·  Operator: Jonathan Brock, Managing Member  ·  plainlydigital.com",
        Inches(0.6),
        Inches(7.05),
        Inches(12),
        Inches(0.3),
        size=10,
        color=RGBColor(0x9C, 0xA3, 0xAF),
    )


def slide_02_thesis():
    s = prs.slides.add_slide(BLANK)
    y = slide_header(
        s, "The thesis", "Three forces are converging. Patet sits exactly in the middle."
    )

    # 3 columns
    col_w = Inches(3.95)
    gap = Inches(0.25)
    col_top = Inches(2.3)
    col_h = Inches(4.4)
    headlines = [
        ("01", "AI quality crossed the line", [
            "Frontier LLMs (Claude Sonnet, GPT-4o) can now answer real personal-finance questions accurately and on-budget.",
            "Inference cost has fallen >50% YoY — viable at sub-$3 MRR.",
            "Coach-style products (Cleo, Origin) have validated demand.",
        ]),
        ("02", "Mint left a vacuum", [
            "Intuit shut Mint down March 2024. ~20M users displaced.",
            "Premium replacements (Monarch, Copilot) cost $13–$15/mo. None of them ship a curriculum.",
            "The free-tier hole is wide open below $5/mo.",
        ]),
        ("03", "Gen Z is unbanked on knowledge", [
            "Only 24% of Gen Z feel confident managing money. 36% find finance “confusing.”",
            "Gen Z gamified-literacy market: $0.31B (2025) → $1.92B (2034), 22.9% CAGR.",
            "No incumbent owns this segment with both bank data + curriculum.",
        ]),
    ]
    x = MARGIN
    for label, head, lines in headlines:
        # background card
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, col_top, col_w, col_h)
        _set_solid(card, GREEN_TINT)
        add_text(
            s,
            label,
            x + Inches(0.3),
            col_top + Inches(0.25),
            Inches(2),
            Inches(0.45),
            size=14,
            bold=True,
            color=GREEN,
        )
        add_text(
            s,
            head,
            x + Inches(0.3),
            col_top + Inches(0.7),
            col_w - Inches(0.6),
            Inches(0.8),
            size=18,
            bold=True,
            color=INK,
        )
        add_bullets(
            s,
            lines,
            x + Inches(0.3),
            col_top + Inches(1.6),
            col_w - Inches(0.6),
            col_h - Inches(1.7),
            size=12,
            color=INK,
        )
        x += col_w + gap

    slide_footer(s, 2, TOTAL_SLIDES)


def slide_03_market():
    s = prs.slides.add_slide(BLANK)
    slide_header(
        s, "Market context", "A real category, an underserved sub-segment, an obvious wedge."
    )

    # left: bullets
    left_x = MARGIN
    left_w = Inches(6.4)
    add_text(
        s,
        "The category is real, the entry tier is empty.",
        left_x,
        Inches(2.2),
        left_w,
        Inches(0.5),
        size=18,
        bold=True,
        color=INK,
    )
    bullets = [
        "Personal-finance software TAM: $1.08B (2022) → $1.59B (2030), 5.1% CAGR. Source: Grand View Research.",
        "Broader personal-finance-apps market projected $207B (2026) → $507B (2030). Source: Business Research Company.",
        "Gen Z gamified-literacy subsegment is the fastest-growing slice: 22.9% CAGR to 2034.",
        "Mint shutdown (Mar 2024) → ~20M users displaced. Monarch absorbed ~1M (April 2026, $850M valuation, $75M Series B).",
        "Rocket Money grew to 10M+ members post-Mint. Copilot Money grew to 100K+ paying.",
        "Median competitor Pro tier price: $12.99/mo. Median annual: ~$95. Patet enters at $2.99 / $6.99.",
    ]
    add_bullets(
        s,
        bullets,
        left_x,
        Inches(2.85),
        left_w,
        Inches(4.3),
        size=13,
    )

    # right: callout card
    right_x = Inches(7.2)
    right_w = SLIDE_W - right_x - MARGIN
    card_top = Inches(2.2)
    card_h = Inches(4.6)
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, right_x, card_top, right_w, card_h
    )
    _set_solid(card, CHARCOAL)
    add_text(
        s,
        "The gap, in one sentence",
        right_x + Inches(0.35),
        card_top + Inches(0.3),
        right_w - Inches(0.7),
        Inches(0.4),
        size=12,
        bold=True,
        color=GREEN,
    )
    add_text(
        s,
        "No competitor in the US market bundles a structured financial-literacy curriculum, live bank data, and a real LLM coach in a single product priced under $7/month.",
        right_x + Inches(0.35),
        card_top + Inches(0.85),
        right_w - Inches(0.7),
        Inches(2.6),
        size=18,
        bold=True,
        color=PAPER,
    )
    add_text(
        s,
        "Cleo has personality but no curriculum. Origin has SEC-regulated AI but no lessons and $13/mo entry. Ramsey has curriculum but no AI or live bank data. YNAB has bank data but no AI, no lessons, and a $15/mo floor.",
        right_x + Inches(0.35),
        card_top + Inches(3.5),
        right_w - Inches(0.7),
        Inches(1.0),
        size=12,
        color=RGBColor(0xC0, 0xC0, 0xC0),
    )

    slide_footer(s, 3, TOTAL_SLIDES)


def slide_04_one_screen():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "The product", "What ships today — one screen.")

    # 4-column metric strip
    metrics = [
        ("18", "modules"),
        ("121", "lessons"),
        ("36+", "screens"),
        ("EN / ES", "languages"),
    ]
    strip_top = Inches(2.2)
    strip_h = Inches(1.4)
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.6)) / 4
    x = MARGIN
    for big, small in metrics:
        card = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, strip_top, col_w, strip_h
        )
        _set_solid(card, GREEN_TINT)
        add_text(
            s,
            big,
            x,
            strip_top + Inches(0.15),
            col_w,
            Inches(0.7),
            size=28,
            bold=True,
            color=GREEN,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            s,
            small,
            x,
            strip_top + Inches(0.85),
            col_w,
            Inches(0.4),
            size=12,
            color=INK,
            align=PP_ALIGN.CENTER,
        )
        x += col_w + Inches(0.2)

    # 3-column body
    col_top = Inches(3.95)
    col_h = Inches(2.8)
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.4)) / 3
    cols = [
        (
            "Curriculum",
            [
                "18 modules · 121 lessons (~2–5 min each)",
                "242 quiz cards, 50-question final assessment",
                "Modules incl. Paycheck, Debt, Saving, Budgeting, Investing, Credit, Student Loans, Mental-Health-of-money",
                "Spanish parity shipped (LATAM reviewer pass pending)",
            ],
        ),
        (
            "AI coach (Glyphe)",
            [
                "Claude Sonnet 4.6 (Pro+) + Haiku 4.5 (Free/Pro)",
                "Knows your 121 lessons + 11 affiliate products + your Plaid spending",
                "4 personas unlocked by progress: Older Sibling → Strategist → Therapist → Hype Friend",
                "Voice in/out · streaming · follow-up chips",
                "Crisis detection + medical/legal/financial guardrails",
            ],
        ),
        (
            "Bank data + tools",
            [
                "Plaid Link (Transactions + Liabilities), encrypted tokens (AES-256-GCM)",
                "Statement upload — PDF / CSV / Excel, AI-extracted transactions",
                "5 calculators, What-If scenarios, Future You, Money Roast, Money IQ, Salary negotiation",
                "Spending-pattern detection → ranked lesson recs",
                "Peer benchmarks (20-user minimum cohort, k-anonymity)",
            ],
        ),
    ]
    x = MARGIN
    for title, lines in cols:
        add_text(
            s,
            title,
            x,
            col_top,
            col_w,
            Inches(0.4),
            size=14,
            bold=True,
            color=INK,
        )
        add_bullets(
            s, lines, x, col_top + Inches(0.45), col_w, col_h - Inches(0.5), size=11
        )
        x += col_w + Inches(0.2)

    slide_footer(s, 4, TOTAL_SLIDES)


def slide_05_moat():
    s = prs.slides.add_slide(BLANK)
    slide_header(
        s,
        "The moat",
        "Three layers, no competitor has all three.",
    )

    # Big diagram: three overlapping nodes, Patet sits at intersection
    diagram_top = Inches(2.3)

    # Three labeled boxes in a triangle
    boxes = [
        (Inches(2.1), diagram_top + Inches(0.0), "Curriculum",
         "18 modules · 121 lessons · cert exam"),
        (Inches(7.6), diagram_top + Inches(0.0), "Live bank data",
         "Plaid + statement upload · 90d history"),
        (Inches(4.85), diagram_top + Inches(2.2), "AI coach",
         "Claude Sonnet / Haiku · spending-aware"),
    ]
    for x, top, head, sub in boxes:
        card = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, top, Inches(3.6), Inches(1.5)
        )
        _set_solid(card, PAPER)
        card.line.color.rgb = GREEN
        card.line.width = Pt(2)
        add_text(
            s,
            head,
            x,
            top + Inches(0.25),
            Inches(3.6),
            Inches(0.5),
            size=18,
            bold=True,
            color=INK,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            s,
            sub,
            x,
            top + Inches(0.8),
            Inches(3.6),
            Inches(0.5),
            size=12,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )

    # Center label — Patet
    center = s.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(5.55),
        diagram_top + Inches(1.5),
        Inches(2.2),
        Inches(1.2),
    )
    _set_solid(center, GREEN)
    add_text(
        s,
        "PATET",
        Inches(5.55),
        diagram_top + Inches(1.85),
        Inches(2.2),
        Inches(0.5),
        size=22,
        bold=True,
        color=PAPER,
        align=PP_ALIGN.CENTER,
    )

    # Caption underneath
    caption = (
        "Ramsey owns curriculum but has no AI and no real-time bank data. "
        "Cleo owns the AI coach but has no curriculum. Monarch / Rocket / "
        "Copilot own bank data but ship no lessons and no LLM coach. "
        "Patet is the only product at the intersection."
    )
    add_text(
        s,
        caption,
        MARGIN,
        Inches(6.1),
        SLIDE_W - 2 * MARGIN,
        Inches(0.8),
        size=13,
        color=INK,
        align=PP_ALIGN.CENTER,
    )

    slide_footer(s, 5, TOTAL_SLIDES)


def slide_06_competitive_matrix():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Competitive landscape", "Where everyone sits — and where the gaps are.")

    headers = [
        "Product",
        "Monthly $",
        "LLM coach",
        "Plaid",
        "Curriculum",
        "Stmt upload",
        "Voice",
    ]
    rows = [
        ("Patet",      "$2.99 / $6.99", "Yes (Sonnet+Haiku)", "Yes",          "Yes (121 lessons)",  "Yes (PDF/CSV/XLS)", "Yes"),
        ("YNAB",       "$14.99",        "No",                 "Yes",          "Workshops only",     "CSV only",          "No"),
        ("Rocket",     "$7–$14",        "No",                 "Yes",          "No",                 "No",                "No"),
        ("Copilot",    "$13",           "Partial (ML)",       "Yes",          "No",                 "CSV only",          "No"),
        ("Monarch",    "$14.99",        "No",                 "Yes",          "No",                 "CSV only",          "No"),
        ("EveryDollar","$17.99",        "No",                 "Premium only", "Ramsey FPU video",   "No",                "No"),
        ("Cleo",       "$5.99–$14.99",  "Yes (multi-model)",  "Yes",          "No",                 "No",                "Limited"),
        ("Origin",     "$12.99",        "Yes (SEC-reg)",      "Yes",          "No",                 "CSV only",          "No"),
    ]

    # Table layout
    table_left = MARGIN
    table_top = Inches(2.25)
    table_w = SLIDE_W - 2 * MARGIN
    col_widths = [
        Inches(1.6),  # Product
        Inches(1.6),  # Price
        Inches(2.1),  # LLM
        Inches(1.3),  # Plaid
        Inches(2.2),  # Curriculum
        Inches(2.0),  # Stmt
        Inches(1.3),  # Voice
    ]
    row_h = Inches(0.5)
    header_h = Inches(0.45)

    # Header row
    x = table_left
    for i, h in enumerate(headers):
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, table_top, col_widths[i], header_h)
        _set_solid(cell, CHARCOAL)
        add_text(
            s,
            h,
            x + Inches(0.1),
            table_top,
            col_widths[i] - Inches(0.2),
            header_h,
            size=11,
            bold=True,
            color=PAPER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        x += col_widths[i]

    # Body rows
    y = table_top + header_h
    for ri, row in enumerate(rows):
        x = table_left
        is_patet = row[0] == "Patet"
        bg = GREEN_TINT if is_patet else (RGBColor(0xF9, 0xFA, 0xFB) if ri % 2 == 0 else PAPER)
        for ci, val in enumerate(row):
            cell = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, col_widths[ci], row_h
            )
            _set_solid(cell, bg)
            color = INK
            bold = is_patet or ci == 0
            add_text(
                s,
                val,
                x + Inches(0.1),
                y,
                col_widths[ci] - Inches(0.2),
                row_h,
                size=10,
                bold=bold,
                color=color,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            x += col_widths[ci]
        y += row_h

    # Caption
    add_text(
        s,
        "Patet is the only row that says Yes across LLM coach, Plaid, curriculum, statement upload, and voice — at the lowest price point. Verified May 2026 via vendor pricing pages.",
        MARGIN,
        Inches(6.7),
        SLIDE_W - 2 * MARGIN,
        Inches(0.5),
        size=11,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    slide_footer(s, 6, TOTAL_SLIDES)


def slide_07_unique():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Unique features", "What we ship that nobody else does.")

    diffs = [
        ("Spending-aware AI coach",
         "Claude Sonnet 4.6 with your real Plaid transactions in-context. Not a chatbot — a coach that can say “you spent $340 on restaurants this month; here’s lesson 4.3.”"),
        ("Statement upload, multi-format",
         "PDF, CSV, and Excel bank statements parsed and AI-categorized. Most competitors offer CSV only — or no upload at all."),
        ("121-lesson curriculum at $2.99",
         "Ramsey charges $79.99/year for curriculum alone. Patet ships 18 modules + 121 lessons in the $2.99/mo Pro tier — and unlocks an AI tutor that knows the curriculum cold."),
        ("Learn-then-unlock coach personas",
         "Four coach personalities unlock at module milestones (Older Sibling default → Strategist @5 → Therapist @10 → Hype Friend @15). Gamifies completion."),
        ("Money Roast + Future You + Money Personality",
         "Emotional engagement layer that competitors don’t ship. Roast gives a 3-tone spending critique. Future You projects your life at 30. Personality maps you to 5 archetypes."),
        ("Voice in / voice out",
         "Native Web Speech on web; expo-speech on mobile. Full voice-mode conversation. Cleo has voice; no other competitor in the matrix does."),
        ("Crisis detection with text normalization",
         "Unicode + homoglyph + l33tspeak normalization on every coach turn. Routes to 988 / Crisis Text Line / 211.org. NIST AI RMF aligned."),
        ("Top-of-funnel viral assessment",
         "Public Financial IQ quiz at /iq (no auth, dynamic OG share image). Free funnel that converts to registered users."),
        ("Spanish curriculum (LATAM)",
         "All 121 lessons + content surfaces translated to LATAM Spanish. Native reviewer pass pending pre-launch. No major US competitor ships Spanish curriculum."),
        ("Strict-freemium tier gating",
         "Plaid + statement upload + advanced calculators + scenarios all sit behind $2.99 Pro. Hard-gated server-side, not trial-and-revoke. Real boundary."),
    ]

    # Two-column grid (5 each)
    col_top = Inches(2.15)
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.4)) / 2
    row_h = Inches(0.95)

    for i, (head, body) in enumerate(diffs):
        col = i % 2
        row = i // 2
        x = MARGIN + col * (col_w + Inches(0.4))
        y = col_top + row * row_h
        # number chip
        num = s.shapes.add_shape(
            MSO_SHAPE.OVAL, x, y + Inches(0.05), Inches(0.32), Inches(0.32)
        )
        _set_solid(num, GREEN)
        add_text(
            s,
            f"{i + 1:02d}",
            x,
            y + Inches(0.08),
            Inches(0.32),
            Inches(0.3),
            size=10,
            bold=True,
            color=PAPER,
            align=PP_ALIGN.CENTER,
        )
        # head
        add_text(
            s,
            head,
            x + Inches(0.45),
            y,
            col_w - Inches(0.45),
            Inches(0.32),
            size=13,
            bold=True,
            color=INK,
        )
        # body
        add_text(
            s,
            body,
            x + Inches(0.45),
            y + Inches(0.32),
            col_w - Inches(0.45),
            Inches(0.6),
            size=10,
            color=MUTED,
        )

    slide_footer(s, 7, TOTAL_SLIDES)


def slide_08_pricing():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Pricing", "Below every competitor's flagship Pro — by 75%.")

    # left: chart
    left_x = MARGIN
    chart_w = Inches(7.0)
    chart_top = Inches(2.3)
    chart_h = Inches(4.4)

    # background
    bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, chart_top, chart_w, chart_h)
    _set_solid(bg, RGBColor(0xF9, 0xFA, 0xFB))

    # bars: [(name, price, color)]
    bars = [
        ("EveryDollar", 17.99, MUTED),
        ("YNAB",        14.99, MUTED),
        ("Monarch",     14.99, MUTED),
        ("Cleo Builder",14.99, MUTED),
        ("Albert",      14.99, MUTED),
        ("Origin",      12.99, MUTED),
        ("Copilot",     13.00, MUTED),
        ("PocketGuard", 12.99, MUTED),
        ("Goodbudget",  10.00, MUTED),
        ("Rocket Money", 9.99, MUTED),
        ("Cleo Plus",    5.99, MUTED),
        ("Patet Pro+",   6.99, GREEN),
        ("Patet Pro",    2.99, GREEN),
    ]
    # Sort: Patet rows pinned at bottom
    competitors = [b for b in bars if not b[0].startswith("Patet")]
    competitors.sort(key=lambda b: -b[1])
    patet_rows = [b for b in bars if b[0].startswith("Patet")]
    ordered = competitors + patet_rows

    max_price = 20.0
    row_h = Inches(0.27)
    row_gap = Inches(0.03)
    chart_inner_top = chart_top + Inches(0.3)
    label_w = Inches(1.6)
    bar_left = left_x + Inches(0.4) + label_w + Inches(0.1)
    bar_max_w = chart_w - (bar_left - left_x) - Inches(1.0)

    for i, (name, price, color) in enumerate(ordered):
        y = chart_inner_top + i * (row_h + row_gap)
        # label
        add_text(
            s,
            name,
            left_x + Inches(0.4),
            y,
            label_w,
            row_h,
            size=10,
            bold=name.startswith("Patet"),
            color=INK if name.startswith("Patet") else MUTED,
            anchor=MSO_ANCHOR.MIDDLE,
            align=PP_ALIGN.RIGHT,
        )
        # bar
        bar_w_emu = int(bar_max_w * (price / max_price))
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left, y + Emu(38100), bar_w_emu, row_h - Emu(76200))
        _set_solid(bar, color)
        # price label
        add_text(
            s,
            f"${price:.2f}",
            bar_left + bar_w_emu + Emu(50800),
            y,
            Inches(0.9),
            row_h,
            size=10,
            bold=name.startswith("Patet"),
            color=GREEN if name.startswith("Patet") else MUTED,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # right: callouts
    right_x = Inches(7.8)
    right_w = SLIDE_W - right_x - MARGIN

    add_text(
        s,
        "What this funds",
        right_x,
        Inches(2.3),
        right_w,
        Inches(0.5),
        size=18,
        bold=True,
        color=INK,
    )
    add_bullets(
        s,
        [
            "Free: 10 coach calls/day on Haiku 4.5. Full curriculum.",
            "Pro $2.99: Plaid, statement upload, all calculators, scenarios, credit tracker, 40 coach calls/day.",
            "Pro+ $6.99: 80 coach calls/day, Claude Sonnet 4.6 (50/mo cap then Haiku fallback), 1,200/mo.",
            "Annual: Pro $24.99 (saves 31%) · Pro+ $57.99 (saves 31%).",
            "$19 one-time IAP: Patet Certified™ credential with LinkedIn add-to-profile.",
            "Boost packs: +50 calls $1.99 · +200 calls $4.99.",
        ],
        right_x,
        Inches(2.95),
        right_w,
        Inches(4.0),
        size=11,
    )

    slide_footer(s, 8, TOTAL_SLIDES)


def slide_09_plaid_core():
    s = prs.slides.add_slide(BLANK)
    slide_header(
        s,
        "Why Plaid is core, not bolt-on",
        "Plaid powers half of Patet's product surface. Removing it would gut Pro.",
    )

    # left: feature list
    add_text(
        s,
        "Plaid-dependent features",
        MARGIN,
        Inches(2.2),
        Inches(6),
        Inches(0.4),
        size=16,
        bold=True,
        color=INK,
    )
    add_bullets(
        s,
        [
            "Bank-account linking via Plaid Link (MFA-gated client + server)",
            "Transaction sync (Transactions product, 90-day history + ongoing webhooks)",
            "Liabilities sync (credit cards, student loans, mortgages, auto loans)",
            "Spending-pattern detection (5 pattern types) → ranked lesson recommendations",
            "Peer benchmarks (anonymized, k≥20 cohort minimum)",
            "Coach context: every AI coach turn has Plaid spending data injected into the system prompt",
            "Spending alert notifications (threshold + category-based, deep-linkable)",
            "Webhook real-time updates (ES256 JWT verified, request_body_sha256 claim verified)",
        ],
        MARGIN,
        Inches(2.7),
        Inches(6),
        Inches(4),
        size=12,
    )

    # right: integration posture
    right_x = Inches(7.0)
    right_w = SLIDE_W - right_x - MARGIN
    card_top = Inches(2.2)
    card_h = Inches(4.6)
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, right_x, card_top, right_w, card_h
    )
    _set_solid(card, CHARCOAL)
    add_text(
        s,
        "Plaid posture",
        right_x + Inches(0.35),
        card_top + Inches(0.25),
        right_w - Inches(0.7),
        Inches(0.45),
        size=12,
        bold=True,
        color=GREEN,
    )
    add_paragraphs(
        s,
        [
            ("Products in use", True),
            ("Transactions, Liabilities · US country code only.", False),
            ("", False),
            ("Security", True),
            ("Plaid access tokens encrypted at rest (AES-256-GCM, dual-write to *_encrypted columns, migration 017).", False),
            ("", False),
            ("Access gate", True),
            ("MFA enforced before /api/plaid/link-token. Email-verified required. 5-min step-up window for Plaid Link replay.", False),
            ("", False),
            ("Webhook verification", True),
            ("ES256 JWT signature, fetched per kid via webhookVerificationKeyGet, JWK→PEM via Node crypto, 5-min replay window.", False),
            ("", False),
            ("Off-boarding", True),
            ("Bank disconnect calls Plaid itemRemove first, then deletes local row. CASCADE on account deletion.", False),
        ],
        right_x + Inches(0.35),
        card_top + Inches(0.7),
        right_w - Inches(0.7),
        card_h - Inches(1),
        size=11,
        color=PAPER,
    )

    slide_footer(s, 9, TOTAL_SLIDES)


def slide_10_security():
    s = prs.slides.add_slide(BLANK)
    slide_header(
        s,
        "Security & compliance",
        "We've already done the work most pre-seed startups skip.",
    )

    items = [
        ("Authentication", [
            "TOTP MFA (AES-256-GCM secret · SHA-256-HMAC recovery codes · 10 codes, single-use)",
            "MFA enforced as a hard gate before Plaid Link (5-min step-up)",
            "Account lockout (DB-persisted, survives Render redeploys)",
            "JWT revocation by jti (migration 016)",
            "httpOnly + sameSite cookies on web; Bearer header on mobile",
        ]),
        ("Data protection", [
            "AES-256-GCM at rest on Plaid tokens, financial PII, MFA secrets",
            "Annual encryption-key rotation procedure documented (InfoSec §7)",
            "Self-serve user data export (GET /api/auth/me/export, JSON, redactions enforced)",
            "Right-to-erasure: DELETE /api/auth/me · FK ON DELETE CASCADE across 30+ child tables",
            "Backups: Neon 7-day window, age out aligned with EDPB 06/2020",
        ]),
        ("Governance & audit", [
            "RBAC (user / support / admin) · system_audit_log table · per-session inventory",
            "NIST AI RMF–aligned audit on every AI coach turn",
            "Crisis detection: Unicode, homoglyph, l33tspeak normalization",
            "Input + output guardrails (prompt injection, role override, medical/legal/financial boundary checks)",
            "1,043 tests across the stack (632 app · 411 server)",
        ]),
        ("Policy library (LLC repo)", [
            "Information Security Policy",
            "Access Control Policy (Plaid-evidence ready)",
            "Data Retention and Disposal Policy (added 2026-05-17 for Plaid IAM Q11)",
            "Privacy Policy live at plainlydigital.com/patet/privacy",
            "Plaid form answers reference packet (PLAID_FORM_ANSWERS.md)",
        ]),
    ]

    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2
    row_h = Inches(2.3)
    starting_top = Inches(2.15)

    for i, (title, lines) in enumerate(items):
        col = i % 2
        row = i // 2
        x = MARGIN + col * (col_w + Inches(0.3))
        y = starting_top + row * row_h
        add_text(
            s,
            title,
            x,
            y,
            col_w,
            Inches(0.4),
            size=13,
            bold=True,
            color=GREEN,
        )
        add_bullets(
            s,
            lines,
            x,
            y + Inches(0.45),
            col_w,
            row_h - Inches(0.5),
            size=10,
        )

    slide_footer(s, 10, TOTAL_SLIDES)


def slide_11_roadmap():
    s = prs.slides.add_slide(BLANK)
    slide_header(
        s,
        "What's next",
        "Six-month roadmap, in order of confidence.",
    )

    items = [
        ("Now → Month 1",
         "Consumer launch",
         "Patet ships strict-freemium to public (live as of 2026-05-16). RevenueCat entitlements + Plaid prod creds + Resend transactional email all wired this week. App Store + Play Store submission queued."),
        ("Month 1 → 3",
         "Affiliate revenue layer",
         "11 affiliate partners across 7 categories (Self, SoFi, Marcus, YNAB, Fidelity, Discover, Ally, Chime, Wealthfront, Lemonade, Schwab). ROAS dashboard live at /api/admin/affiliate/roas. Per-partner conversion webhooks live."),
        ("Month 2 → 6",
         "Patet Certified™",
         "50-question final assessment shipped. $19 one-time IAP via RevenueCat. LinkedIn add-to-profile URL helper. Vanity slug picker. Public verify endpoint. Credential as a B2B credential layer for community-college + NACFC coach licensing."),
        ("Month 3 → 6",
         "LATAM launch",
         "Spanish lesson content shipped (machine-translated via Sonnet, native reviewer pass pending). 100% i18n parity test-enforced. Expansion candidate: Mexico, Colombia, Argentina."),
        ("Month 4 → 6",
         "B2B curriculum licensing",
         "NACFC coach pilot (financial coaches license Patet curriculum for their clients). Community-college pilot. Employer benefits pilot via SHRM expo channel."),
        ("Month 6+",
         "GCP migration",
         "Move backend off Render onto Cloud Run + Cloud SQL. Cloud Build CI on push. Already migrated other Plainly Digital products; Patet is on the queue."),
    ]

    top = Inches(2.15)
    row_h = Inches(0.78)
    col_when_w = Inches(2.0)
    col_what_w = Inches(2.8)
    col_body_w = SLIDE_W - 2 * MARGIN - col_when_w - col_what_w - Inches(0.3)

    for i, (when, what, body) in enumerate(items):
        y = top + i * row_h
        # when chip
        chip = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            MARGIN,
            y + Inches(0.1),
            col_when_w,
            Inches(0.45),
        )
        _set_solid(chip, GREEN_TINT)
        add_text(
            s,
            when,
            MARGIN,
            y + Inches(0.13),
            col_when_w,
            Inches(0.4),
            size=11,
            bold=True,
            color=GREEN,
            align=PP_ALIGN.CENTER,
        )
        # what
        add_text(
            s,
            what,
            MARGIN + col_when_w + Inches(0.15),
            y + Inches(0.1),
            col_what_w,
            Inches(0.45),
            size=15,
            bold=True,
            color=INK,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # body
        add_text(
            s,
            body,
            MARGIN + col_when_w + col_what_w + Inches(0.3),
            y + Inches(0.1),
            col_body_w,
            row_h - Inches(0.2),
            size=10,
            color=MUTED,
        )

    slide_footer(s, 11, TOTAL_SLIDES)


def slide_12_close():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, CHARCOAL)

    # accent block
    add_rect(s, Inches(0.6), Inches(2.4), Inches(0.18), Inches(2.4), GREEN)

    add_text(
        s,
        "Thanks for partnering with us.",
        Inches(1.0),
        Inches(2.4),
        Inches(11.3),
        Inches(0.9),
        size=36,
        bold=True,
        color=PAPER,
    )
    add_text(
        s,
        "Plaid is the load-bearing partner in the most differentiated half of our product. We are building this slowly and carefully — security, governance, and customer trust first. Happy to walk through any specific area in more depth.",
        Inches(1.0),
        Inches(3.4),
        Inches(11.3),
        Inches(2.0),
        size=15,
        color=RGBColor(0xD1, 0xD5, 0xDB),
    )

    # contact card
    contact_x = Inches(1.0)
    contact_y = Inches(5.4)
    contact_w = Inches(11.3)
    card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, contact_x, contact_y, contact_w, Inches(1.4)
    )
    _set_solid(card, RGBColor(0x27, 0x27, 0x2A))
    add_paragraphs(
        s,
        [
            ("Plainly Digital LLC  ·  DBA Patet", True),
            ("Operator: Jonathan Brock, Managing Member", False),
            ("apps@plainlydigital.com  ·  support@plainlydigital.com  ·  plainlydigital.com/patet", False),
        ],
        contact_x + Inches(0.4),
        contact_y + Inches(0.2),
        contact_w - Inches(0.8),
        Inches(1),
        size=13,
        color=PAPER,
    )

    add_text(
        s,
        "12 / 12",
        SLIDE_W - MARGIN - Inches(1),
        Inches(7.05),
        Inches(1),
        Inches(0.3),
        size=9,
        color=RGBColor(0x6B, 0x72, 0x80),
        align=PP_ALIGN.RIGHT,
    )


# ─── Build ────────────────────────────────────────────────────────────────────

slide_01_title()
slide_02_thesis()
slide_03_market()
slide_04_one_screen()
slide_05_moat()
slide_06_competitive_matrix()
slide_07_unique()
slide_08_pricing()
slide_09_plaid_core()
slide_10_security()
slide_11_roadmap()
slide_12_close()

prs.save(OUTPUT)
print(f"Wrote {OUTPUT}")
